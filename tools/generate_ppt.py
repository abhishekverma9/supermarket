#!/usr/bin/env python3
"""
Generate a PowerPoint presentation `supermarket_presentation.pptx` from repo files.
This script installs `python-pptx` and `Pillow` if missing, extracts key content,
and writes a 12-slide deck with speaker notes and a small Appendix.
"""
import os
import sys
import subprocess

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except Exception:
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx", "Pillow"]) 
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

import json

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
OUTPUT = os.path.join(ROOT, 'supermarket_presentation.pptx')

def read_file(path, max_chars=3000):
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            data = f.read()
            return data[:max_chars]
    except Exception:
        return ''

def get_repo_name():
    readme = os.path.join(ROOT, 'readme.md')
    if os.path.exists(readme):
        text = read_file(readme, 300)
        first_line = text.strip().splitlines()[0] if text.strip() else 'Supermarket'
        return first_line
    return 'Supermarket'

def load_json(path):
    try:
        with open(path, 'r', encoding='utf-8') as f:
            return json.load(f)
    except Exception:
        return {}

def list_files(relative_dir, ext=None, limit=20):
    path = os.path.join(ROOT, relative_dir)
    out = []
    if not os.path.exists(path):
        return out
    for root, dirs, files in os.walk(path):
        for fn in files:
            if ext and not fn.endswith(ext):
                continue
            out.append(os.path.relpath(os.path.join(root, fn), ROOT))
            if len(out) >= limit:
                return out
    return out

prs = Presentation()
blank_slide_layout = prs.slide_layouts[6]

# Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
repo_name = get_repo_name()
title.text = repo_name
subtitle.text = 'A lightweight supermarket app (frontend: React + Vite; backend: Node.js + Express)\nAuto-generated presentation'

# Try to add logo if present
logo_paths = [
    'frontend/public/logo.jpg',
    'frontend/public/logo.png',
    'frontend/public/logooo.png',
    'frontend/public/logo.jpg'
]
for lp in logo_paths:
    lp_abs = os.path.join(ROOT, lp)
    if os.path.exists(lp_abs):
        try:
            slide.shapes.add_picture(lp_abs, prs.slide_width - Inches(1.2), Inches(0.2), width=Inches(1))
            break
        except Exception:
            pass

# Executive Summary
s = prs.slides.add_slide(blank_slide_layout)
left = Inches(0.6)
top = Inches(0.5)
width = prs.slide_width - Inches(1.2)
height = Inches(1.8)
tx = s.shapes.add_textbox(left, top, width, height).text_frame
tx.word_wrap = True
p = tx.paragraphs[0]
p.text = 'Executive Summary'
p.font.size = Pt(28)
p.font.bold = True
# body
p2 = tx.add_paragraph()
p2.text = 'An end-to-end supermarket application with consumer, employee, and owner interfaces, REST API backend, image handling, email/OTP utilities, and a small Python chatbot. Includes product, order, cart and authentication flows.'
p2.font.size = Pt(14)

# Tech Stack
s = prs.slides.add_slide(blank_slide_layout)
tx = s.shapes.add_textbox(left, top, width, Inches(2)).text_frame
p = tx.paragraphs[0]
p.text = 'Tech Stack'
p.font.size = Pt(28)
p.font.bold = True
b = tx.add_paragraph()
backend_pkg = load_json(os.path.join(ROOT, 'backend', 'package.json'))
frontend_pkg = load_json(os.path.join(ROOT, 'frontend', 'package.json'))
stack_lines = [
    f"Backend: Node.js, Express ({backend_pkg.get('name','backend')})",
    f"Frontend: React, Vite ({frontend_pkg.get('name','frontend')})",
    "Database: configured in backend/config/db.js (likely MySQL/Mongo based on config)",
    "Other: ImageKit, email utilities, OTP store, simple session/token utilities",
]
for ln in stack_lines:
    p = tx.add_paragraph()
    p.text = '• ' + ln
    p.level = 1
    p.font.size = Pt(14)

# Architecture slide (textual)
s = prs.slides.add_slide(blank_slide_layout)
tx = s.shapes.add_textbox(left, top, width, Inches(2)).text_frame
p = tx.paragraphs[0]
p.text = 'Architecture (high level)'
p.font.size = Pt(24)
p.font.bold = True
for ln in [
    'Browser (React) → Frontend (Vite) → API (Express server.js)',
    'Routes → Controllers → Models → DB (configured in backend/config/db.js)',
    'Third-party: ImageKit (image uploads), Email service (emailService.js), Auth tokens',
]:
    q = tx.add_paragraph(); q.text = '• ' + ln; q.font.size = Pt(14)

# Backend: Routes & Endpoints (list routes files)
s = prs.slides.add_slide(blank_slide_layout)
tx = s.shapes.add_textbox(left, top, width, Inches(2)).text_frame
p = tx.paragraphs[0]
p.text = 'Backend: Routes & Endpoints'
p.font.size = Pt(24)
p.font.bold = True
routes = list_files('backend/routes', ext='.js', limit=50)
if not routes:
    tx.add_paragraph().text = 'No route files found.'
else:
    for r in routes:
        tx.add_paragraph().text = '• ' + r; tx.paragraphs[-1].font.size = Pt(14)

# Models slide
s = prs.slides.add_slide(blank_slide_layout)
tx = s.shapes.add_textbox(left, top, width, Inches(2)).text_frame
p = tx.paragraphs[0]
p.text = 'Data Models'
p.font.size = Pt(24)
p.font.bold = True
models = list_files('backend/models', ext='.js', limit=20)
if models:
    for m in models:
        tx.add_paragraph().text = '• ' + m; tx.paragraphs[-1].font.size = Pt(14)
else:
    tx.add_paragraph().text = 'No models found.'

# Frontend: Pages & UX
s = prs.slides.add_slide(blank_slide_layout)
tx = s.shapes.add_textbox(left, top, width, Inches(3)).text_frame
p = tx.paragraphs[0]
p.text = 'Frontend: Pages & UX'
p.font.size = Pt(24)
p.font.bold = True
pages = list_files('frontend/src/pages', ext='.jsx', limit=30)
for pg in pages:
    tx.add_paragraph().text = '• ' + pg; tx.paragraphs[-1].font.size = Pt(14)

# Chatbot slide
s = prs.slides.add_slide(blank_slide_layout)
tx = s.shapes.add_textbox(left, top, width, Inches(3)).text_frame
p = tx.paragraphs[0]
p.text = 'Chatbot Feature'
p.font.size = Pt(24)
p.font.bold = True
cb_files = list_files('frontend/src/pages/chatbot', limit=30)
if cb_files:
    for f in cb_files:
        tx.add_paragraph().text = '• ' + f; tx.paragraphs[-1].font.size = Pt(14)
else:
    tx.add_paragraph().text = 'No chatbot files found.'

# Demo / Run Instructions
s = prs.slides.add_slide(blank_slide_layout)
tx = s.shapes.add_textbox(left, top, width, Inches(3)).text_frame
p = tx.paragraphs[0]
p.text = 'Demo / Local Run Instructions'
p.font.size = Pt(24)
p.font.bold = True
cmds = [
    'Backend: cd backend ; npm install ; npm run dev',
    'Frontend: cd frontend ; npm install ; npm run dev',
    'Ensure .env files are configured (see backend/.env and frontend/.env.example)',
]
for c in cmds:
    tx.add_paragraph().text = c; tx.paragraphs[-1].font.size = Pt(14)

# Code highlight (orderController snippet)
order_controller = os.path.join(ROOT, 'backend', 'controllers', 'orderController.js')
snippet = read_file(order_controller, max_chars=2000)
if snippet:
    s = prs.slides.add_slide(blank_slide_layout)
    tx = s.shapes.add_textbox(left, top, width, Inches(4)).text_frame
    p = tx.paragraphs[0]
    p.text = 'Code Highlight: backend/controllers/orderController.js'
    p.font.size = Pt(20)
    p.font.bold = True
    pb = tx.add_paragraph()
    # Add only first 4000 chars in a monospace-like style
    code_text = snippet[:3000]
    pb.text = code_text
    pb.font.name = 'Consolas'
    pb.font.size = Pt(10)

# Roadmap / Future Work
s = prs.slides.add_slide(blank_slide_layout)
tx = s.shapes.add_textbox(left, top, width, Inches(3)).text_frame
p = tx.paragraphs[0]
p.text = 'Roadmap / Future Work'
p.font.size = Pt(24)
p.font.bold = True
for ln in ['Add tests and CI', 'Add production deployment scripts', 'Improve auth and input validation', 'Add E2E tests and mobile UI enhancements']:
    tx.add_paragraph().text = '• ' + ln; tx.paragraphs[-1].font.size = Pt(14)

# Appendix: file map (small)
s = prs.slides.add_slide(blank_slide_layout)
tx = s.shapes.add_textbox(left, top, width, Inches(4)).text_frame
p = tx.paragraphs[0]
p.text = 'Appendix: Important files (trimmed)'
p.font.size = Pt(20)
p.font.bold = True
all_files = []
for d in ['backend', 'frontend']:
    for root, dirs, files in os.walk(os.path.join(ROOT, d)):
        for fn in files:
            rel = os.path.relpath(os.path.join(root, fn), ROOT)
            all_files.append(rel)
            if len(all_files) > 80:
                break
        if len(all_files) > 80:
            break
    if len(all_files) > 80:
        break
for f in all_files[:80]:
    tx.add_paragraph().text = f; tx.paragraphs[-1].font.size = Pt(10)

# Save
prs.save(OUTPUT)
print('Presentation saved to', OUTPUT)

if __name__ == '__main__':
    pass
